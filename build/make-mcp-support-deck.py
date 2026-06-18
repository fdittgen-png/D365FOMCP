#!/usr/bin/env python
"""Generate a short PowerPoint deck summarising the locally-available D365FO MCP
services and how they improve the Claude experience for support.

Run: python build/make-mcp-support-deck.py
Output: docs/D365FO-MCP-for-Support.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette (Trelleborg-ish blue / slate) ─────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x21, 0x40)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)
ACCENT = RGBColor(0x16, 0xA0, 0x85)
SLATE  = RGBColor(0x45, 0x52, 0x60)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x8A, 0x97, 0xA5)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (text,size,color,bold,italic)."""
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = "Segoe UI"
            if rest and rest[0]:
                r.font.italic = True
    return tb


def header(s, kicker, title):
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), SW, Pt(3), ACCENT)
    txt(s, Inches(0.6), Inches(0.12), Inches(12), Inches(0.35),
        [[(kicker, 12, ACCENT, True)]])
    txt(s, Inches(0.6), Inches(0.40), Inches(12.1), Inches(0.7),
        [[(title, 26, WHITE, True)]])


def bullets(s, x, y, w, h, items, size=14, gap=8, color=SLATE):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            paras.append([("•  ", size, ACCENT, True),
                          (lead, size, NAVY, True),
                          (rest, size, color, False)])
        else:
            paras.append([("•  ", size, ACCENT, True), (it, size, color, False)])
    txt(s, x, y, w, h, paras, space_after=gap, line_spacing=1.05)


# ── 1. Title ──────────────────────────────────────────────────────────────────
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.7), SW, Pt(3), ACCENT)
txt(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.6),
    [[("D365FO MCP Services", 46, WHITE, True)],
     [("Grounding Claude in real D365 metadata — for faster, accurate support", 22, RGBColor(0xBF,0xD6,0xEA), False)]],
    space_after=14)
txt(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.2),
    [[("54 AI-consumable tools · 4 local services · read-only · grounded in your environment", 15, ACCENT, True)],
     [("Internal enablement overview · 2026", 12, GREY, False)]], space_after=8)

# ── 2. The problem ────────────────────────────────────────────────────────────
s = slide()
header(s, "WHY", "A general LLM alone is risky for D365 support")
bullets(s, Inches(0.7), Inches(1.5), Inches(7.6), Inches(5.4), [
    ("Hallucinations. ", "Generic models invent table names, fields, enum values and role names that don't exist in your build."),
    ("No environment truth. ", "They can't see your AOT metadata, cross-references, security config or customizations (TBG_/TOC_)."),
    ("Version drift. ", "Field renames and customizations differ per environment — generic answers go stale."),
    ("Slow triage. ", "Engineers hand-trace callers, permissions and repro steps across tools."),
], size=15, gap=14)
rect(s, Inches(8.7), Inches(1.6), Inches(3.9), Inches(4.9), LIGHT)
rect(s, Inches(8.7), Inches(1.6), Inches(3.9), Pt(3), BLUE)
txt(s, Inches(8.95), Inches(1.85), Inches(3.5), Inches(4.5),
    [[("The fix", 16, BLUE, True)],
     [("Give Claude live, read-only access to the same metadata an engineer would open in Visual Studio — via the Model Context Protocol (MCP).", 14, SLATE, False)],
     [("Every answer is traceable to a real object in the snapshot.", 14, NAVY, True)]],
    space_after=12, line_spacing=1.1)

# ── 3. What is the platform ───────────────────────────────────────────────────
s = slide()
header(s, "WHAT", "One MCP platform, four services, 54 tools")
cards = [
    ("Knowledge Base", "17 tools", "Tables, fields, EDTs, enums,\nclasses, methods, labels", BLUE),
    ("Cross-Reference", "16 tools", "Callers, usages, extensions,\nevent handlers, impact", ACCENT),
    ("Security", "19 tools", "Roles, duties, privileges,\neffective access, SoD, licence", RGBColor(0xB5,0x4D,0x2E)),
    ("Task Recorder", "2 tools", "Recordings → enriched\nrepro & process docs", RGBColor(0x6B,0x4B,0xA3)),
]
x = Inches(0.55); w = Inches(3.0); gap = Inches(0.13)
for (name, count, desc, col) in cards:
    rect(s, x, Inches(1.7), w, Inches(3.4), LIGHT)
    rect(s, x, Inches(1.7), w, Inches(0.85), col)
    txt(s, x, Inches(1.78), w, Inches(0.8),
        [[(name, 15, WHITE, True)], [(count, 12, WHITE, False)]],
        align=PP_ALIGN.CENTER, space_after=2)
    txt(s, x+Inches(0.15), Inches(2.8), w-Inches(0.3), Inches(2.1),
        [[(line, 13, SLATE, False)] for line in desc.split("\n")],
        align=PP_ALIGN.CENTER, space_after=4, line_spacing=1.1)
    x = x + w + gap
txt(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.6),
    [[("Local & safe: ", 14, NAVY, True),
      ("each service runs as a local stdio MCP server over a read-only SQLite snapshot. Read-only, idempotent, no open-world calls — hosts can auto-approve.", 14, SLATE, False)],
     [("Plus a complementary RAG service ", 13, NAVY, True),
      ("for grounded Q&A over D365 documentation.", 13, SLATE, False)]],
    space_after=10, line_spacing=1.1)

# ── 4. KB ─────────────────────────────────────────────────────────────────────
s = slide()
header(s, "SERVICE 1 · KNOWLEDGE BASE (17 tools)", "Answer 'what is this object?' from real metadata")
bullets(s, Inches(0.7), Inches(1.5), Inches(7.7), Inches(5.4), [
    ("Look it up. ", "lookup_table, get_enum, get_class_methods, get_method_source, resolve_label — exact field/enum/label values, no guessing."),
    ("Explore relations. ", "find_referencing_tables, get_join_keys, get_entity_sources, graph_traverse."),
    ("Catch drift. ", "field_renames and check_field_exists confirm a field is real in this build."),
    ("Anti-hallucination. ", "hallucination_check validates a claimed object before you act on it."),
], size=14, gap=12)
rect(s, Inches(8.7), Inches(1.6), Inches(3.9), Inches(4.9), NAVY)
txt(s, Inches(8.95), Inches(1.85), Inches(3.5), Inches(4.5),
    [[("Support win", 15, ACCENT, True)],
     [("“Is CustTable.Blocked an enum?”", 14, WHITE, True)],
     [("Claude returns the real EDT/enum, its values and label — instead of a confident guess.", 13, RGBColor(0xBF,0xD6,0xEA), False)]],
    space_after=12, line_spacing=1.15)

# ── 5. XRef ───────────────────────────────────────────────────────────────────
s = slide()
header(s, "SERVICE 2 · CROSS-REFERENCE (16 tools)", "Trace 'what touches this?' across the codebase")
bullets(s, Inches(0.7), Inches(1.5), Inches(7.7), Inches(5.4), [
    ("Who calls what. ", "find_references, find_usages, find_method_callers, method_references, find_field_usages."),
    ("Customization-aware. ", "find_extensions and find_event_handlers surface ISV / TBG_ / TOC_ overlayer behaviour."),
    ("Structure. ", "class_hierarchy, interface_implementors, module_objects, cross_module_deps."),
    ("Blast radius. ", "impact_analysis answers “what breaks if we change X?” before a fix ships."),
], size=14, gap=12)
rect(s, Inches(8.7), Inches(1.6), Inches(3.9), Inches(4.9), NAVY)
txt(s, Inches(8.95), Inches(1.85), Inches(3.5), Inches(4.5),
    [[("Support win", 15, ACCENT, True)],
     [("“Why did posting change?”", 14, WHITE, True)],
     [("Claude lists every caller and extension of the method — with line references you can click straight to.", 13, RGBColor(0xBF,0xD6,0xEA), False)]],
    space_after=12, line_spacing=1.15)

# ── 6. Security ───────────────────────────────────────────────────────────────
s = slide()
header(s, "SERVICE 3 · SECURITY (19 tools)", "Explain any access decision — end to end")
bullets(s, Inches(0.7), Inches(1.5), Inches(7.7), Inches(5.4), [
    ("Who & what. ", "lookup_user/role/duty/privilege, find_users_by_role, company_users, role_hierarchy."),
    ("Effective access. ", "effective_permissions & permission_trace walk role→duty→privilege→entry point with Deny-wins."),
    ("“Can user X do Y?” ", "object_access shows grant/deny incl. Invoke; what_if simulates a role change."),
    ("Governance. ", "sod_check (now from live D365 SoD rules), licence_assessment, compare_roles."),
], size=14, gap=12)
rect(s, Inches(8.7), Inches(1.6), Inches(3.9), Inches(4.9), NAVY)
txt(s, Inches(8.95), Inches(1.85), Inches(3.5), Inches(4.5),
    [[("Support win", 15, ACCENT, True)],
     [("“Why is the button greyed out?”", 14, WHITE, True)],
     [("Claude traces the exact role/duty granting or denying the menu item — no more screen-share spelunking.", 13, RGBColor(0xBF,0xD6,0xEA), False)]],
    space_after=12, line_spacing=1.15)

# ── 7. Task Recorder ──────────────────────────────────────────────────────────
s = slide()
header(s, "SERVICE 4 · TASK RECORDER (2 tools)", "Turn a recording into a repro you can act on")
bullets(s, Inches(0.7), Inches(1.5), Inches(7.7), Inches(5.4), [
    ("taskrecorder_to_markdown. ", "Parses a server .axtr into clean, readable steps."),
    ("taskrecorder_to_document. ", "Builds an enriched MHTML web-archive: client recording + screenshots correlated to each server action."),
    ("Auto-enriched. ", "Each step is annotated with KB metadata and Security context (which object, which permission)."),
    ("Reproduce faster. ", "A customer recording becomes a precise, navigable repro for the engineer."),
], size=14, gap=12)
rect(s, Inches(8.7), Inches(1.6), Inches(3.9), Inches(4.9), NAVY)
txt(s, Inches(8.95), Inches(1.85), Inches(3.5), Inches(4.5),
    [[("Support win", 15, ACCENT, True)],
     [("“It fails when I post…”", 14, WHITE, True)],
     [("The .axtr becomes a step-by-step document with the exact forms, menu items and required permissions.", 13, RGBColor(0xBF,0xD6,0xEA), False)]],
    space_after=12, line_spacing=1.15)

# ── 8. Support workflows ──────────────────────────────────────────────────────
s = slide()
header(s, "IN PRACTICE", "How a support engineer uses it")
flow = [
    ("Scope", "Frame the ticket; identify objects & users involved", BLUE),
    ("Reproduce", "Turn the Task Recording into an enriched repro doc", ACCENT),
    ("Diagnose", "KB + XRef trace the behaviour & blast radius", RGBColor(0xB5,0x4D,0x2E)),
    ("Access", "Security explains who can/can't do what, and why", RGBColor(0x6B,0x4B,0xA3)),
]
x = Inches(0.55); w = Inches(3.0); gap = Inches(0.13)
for i, (name, desc, col) in enumerate(flow):
    rect(s, x, Inches(1.8), w, Inches(2.4), LIGHT)
    rect(s, x, Inches(1.8), w, Inches(0.7), col)
    txt(s, x, Inches(1.9), w, Inches(0.5), [[(f"{i+1}. {name}", 15, WHITE, True)]],
        align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.15), Inches(2.65), w-Inches(0.3), Inches(1.4),
        [[(desc, 13, SLATE, False)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
    x = x + w + gap
txt(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(2.4),
    [[("Packaged as one-click skills: ", 14, NAVY, True),
      ("support-scope, support-reproduce, support-diagnose and biz-access-check orchestrate the tools above so engineers don't memorise tool names.", 14, SLATE, False)],
     [("Every claim is grounded: ", 14, NAVY, True),
      ("answers cite the real object and snapshot date, so they're auditable, not improvised.", 14, SLATE, False)]],
    space_after=12, line_spacing=1.12)

# ── 9. Impact / close ─────────────────────────────────────────────────────────
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(1.15), SW, Pt(3), ACCENT)
txt(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.7),
    [[("Why it matters for support", 28, WHITE, True)]])
gains = [
    ("Accurate", "Grounded in real metadata — far fewer hallucinations"),
    ("Faster", "Minutes, not hours, to trace callers / permissions / repro"),
    ("Auditable", "Every answer traces to an object + snapshot date"),
    ("Safe", "Read-only snapshots; no consumer/vendor data, no writes"),
]
x = Inches(0.6); w = Inches(2.95); gap = Inches(0.16)
for (name, desc, ) in [(g[0], g[1]) for g in gains]:
    rect(s, x, Inches(1.7), w, Inches(2.7), RGBColor(0x12,0x2F,0x52))
    rect(s, x, Inches(1.7), w, Pt(3), ACCENT)
    txt(s, x+Inches(0.2), Inches(1.95), w-Inches(0.4), Inches(2.3),
        [[(name, 18, ACCENT, True)], [(desc, 13, RGBColor(0xCF,0xDD,0xEA), False)]],
        space_after=10, line_spacing=1.15)
    x = x + w + gap
txt(s, Inches(0.7), Inches(4.9), Inches(12), Inches(2.2),
    [[("From “confident guesses” to “traceable answers.”", 20, WHITE, True)],
     [("4 services · 54 tools · running locally today — extensible as the D365 estate grows.", 14, GREY, False)]],
    space_after=12)

import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   "docs", "D365FO-MCP-for-Support.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
